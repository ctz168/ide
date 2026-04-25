"""
PhoneIDE - File management API routes.
"""

import json
import os
import re
import fnmatch
from pathlib import Path
from datetime import datetime
from flask import Blueprint, jsonify, request, send_file, Response
from utils import (
    handle_error, load_config, save_config, WORKSPACE,
    get_icon_for_file, get_file_type,
)

bp = Blueprint('files', __name__)


@bp.route('/api/files/list', methods=['GET'])
@handle_error
def list_files():
    path = request.args.get('path', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)
    project = config.get('project', None)

    target = os.path.join(base, path) if path else base
    target = os.path.realpath(target)

    # Security: must be under workspace
    real_base = os.path.realpath(base)
    if not target.startswith(real_base):
        return jsonify({'error': 'Access denied'}), 403

    # Project boundary enforcement:
    # When a project is open, the file tree is confined to the project directory.
    # If the requested path is above the project root, redirect to the project root.
    if project:
        project_dir = os.path.realpath(os.path.join(base, project))
        if os.path.isdir(project_dir) and not target.startswith(project_dir):
            # User tried to navigate above the project — redirect to project root
            target = project_dir
            path = project

    # Auto-create workspace root if it doesn't exist
    if not os.path.exists(target) and target == os.path.realpath(base):
        try:
            os.makedirs(target, exist_ok=True)
        except OSError:
            return jsonify({'error': f'Cannot create workspace directory: {target}'}), 500

    if not os.path.exists(target):
        return jsonify({'error': 'Path not found'}), 404

    items = []
    if os.path.isdir(target):
        try:
            for entry in sorted(os.listdir(target)):
                full = os.path.join(target, entry)
                try:
                    st = os.stat(full)
                    is_dir = os.path.isdir(full)
                    items.append({
                        'name': entry,
                        'path': os.path.relpath(full, base).replace(os.sep, '/'),
                        'is_dir': is_dir,
                        'size': st.st_size if not is_dir else 0,
                        'modified': datetime.fromtimestamp(st.st_mtime).isoformat(),
                        'icon': get_icon_for_file(entry),
                    })
                except (PermissionError, OSError):
                    pass
        except PermissionError:
            return jsonify({'error': 'Permission denied'}), 403
    else:
        items.append({
            'name': os.path.basename(target),
            'path': os.path.relpath(target, base),
            'is_dir': False,
            'size': os.path.getsize(target),
            'modified': datetime.fromtimestamp(os.path.getmtime(target)).isoformat(),
            'icon': get_icon_for_file(os.path.basename(target)),
        })

    return jsonify({'items': items, 'path': path, 'base': base, 'project': project})


@bp.route('/api/files/read', methods=['GET'])
@handle_error
def read_file():
    path = request.args.get('path', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    target = os.path.realpath(os.path.join(base, path))

    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.isfile(target):
        return jsonify({'error': 'File not found'}), 404

    # Limit file size (10MB)
    size = os.path.getsize(target)
    if size > 10 * 1024 * 1024:
        return jsonify({'error': 'File too large (>10MB)', 'size': size}), 413

    try:
        # Try to detect encoding
        with open(target, 'rb') as f:
            raw = f.read()

        encodings = ['utf-8', 'utf-8-sig', 'gbk', 'gb2312', 'latin-1']
        content = None
        used_encoding = 'utf-8'
        for enc in encodings:
            try:
                content = raw.decode(enc)
                used_encoding = enc
                break
            except (UnicodeDecodeError, LookupError):
                continue

        if content is None:
            content = raw.decode('utf-8', errors='replace')
            used_encoding = 'utf-8'

        return jsonify({
            'content': content,
            'path': path,
            'encoding': used_encoding,
            'type': get_file_type(os.path.basename(target)),
            'size': size,
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# MIME type mapping for preview
_PREVIEW_MIME_TYPES = {
    '.html': 'text/html; charset=utf-8',
    '.htm': 'text/html; charset=utf-8',
    '.md': 'text/markdown; charset=utf-8',
    '.markdown': 'text/markdown; charset=utf-8',
    '.css': 'text/css; charset=utf-8',
    '.js': 'application/javascript; charset=utf-8',
    '.json': 'application/json; charset=utf-8',
    '.xml': 'application/xml; charset=utf-8',
    '.svg': 'image/svg+xml; charset=utf-8',
    '.txt': 'text/plain; charset=utf-8',
    '.py': 'text/plain; charset=utf-8',
    '.png': 'image/png',
    '.jpg': 'image/jpeg',
    '.jpeg': 'image/jpeg',
    '.gif': 'image/gif',
    '.ico': 'image/x-icon',
    '.webp': 'image/webp',
    '.pdf': 'application/pdf',
    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
}


@bp.route('/api/files/preview', methods=['GET'])
@handle_error
def preview_file():
    """Serve a local file for browser preview (HTML, MD, images, etc.).
    This route returns raw file content with proper Content-Type so that
    the browser's iframe can render it correctly."""
    path = request.args.get('path', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    target = os.path.realpath(os.path.join(base, path))

    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.isfile(target):
        return jsonify({'error': 'File not found'}), 404

    # Determine MIME type from extension
    ext = os.path.splitext(target)[1].lower()
    mime_type = _PREVIEW_MIME_TYPES.get(ext, 'application/octet-stream')

    # For Markdown files, convert to HTML before serving
    if ext in ('.md', '.markdown'):
        try:
            with open(target, 'r', encoding='utf-8', errors='replace') as f:
                md_content = f.read()
            # Safely encode markdown content as JSON string to prevent XSS
            md_json = json.dumps(md_content)
            html_content = f'''<!DOCTYPE html>
<html><head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<link rel="stylesheet" href="/vendor/katex/katex.min.css">
<link rel="stylesheet" href="/vendor/highlightjs/github.min.css" id="hljs-light">
<link rel="stylesheet" href="/vendor/highlightjs/github-dark.min.css" id="hljs-dark" disabled>
<style>
:root {{ --bg: #fff; --fg: #24292f; --heading: #1a1a2e; --link: #0969da; --code-bg: #f6f8fa;
         --border: #d0d7de; --blockquote: #656d76; --blockquote-bg: #f6f8fa; }}
@media (prefers-color-scheme: dark) {{
  :root {{ --bg: #0d1117; --fg: #c9d1d9; --heading: #f0f6fc; --link: #58a6ff; --code-bg: #161b22;
           --border: #30363d; --blockquote: #8b949e; --blockquote-bg: #161b22; }}
  #hljs-light {{ disabled: true; }}
  #hljs-dark {{ disabled: false; }}
}}
body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
       max-width: 900px; margin: 0 auto; padding: 20px; line-height: 1.7;
       color: var(--fg); background: var(--bg); }}
h1 {{ font-size: 2em; font-weight: 700; margin-top: 1.2em; margin-bottom: 0.4em;
      border-bottom: 2px solid var(--border); padding-bottom: 0.3em; color: var(--heading); }}
h2 {{ font-size: 1.5em; font-weight: 600; margin-top: 1.1em; margin-bottom: 0.35em;
      border-bottom: 1px solid var(--border); padding-bottom: 0.25em; color: var(--heading); }}
h3 {{ font-size: 1.25em; font-weight: 600; margin-top: 1em; margin-bottom: 0.3em; color: var(--heading); }}
h4 {{ font-size: 1.1em; font-weight: 600; margin-top: 0.9em; margin-bottom: 0.25em; color: var(--heading); }}
h5,h6 {{ font-size: 1em; font-weight: 600; margin-top: 0.8em; margin-bottom: 0.2em; color: var(--heading); }}
code {{ background: var(--code-bg); padding: 2px 6px; border-radius: 3px; font-size: 0.9em;
        font-family: "SFMono-Regular", Consolas, "Liberation Mono", Menlo, monospace; }}
pre {{ background: var(--code-bg); padding: 14px; border-radius: 8px; overflow-x: auto;
       border: 1px solid var(--border); margin: 1em 0; }}
pre code {{ background: none; padding: 0; font-size: 13px; line-height: 1.5; }}
blockquote {{ border-left: 4px solid var(--link); margin: 1em 0; padding: 8px 16px;
              color: var(--blockquote); background: var(--blockquote-bg); border-radius: 0 6px 6px 0; }}
table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
th,td {{ border: 1px solid var(--border); padding: 8px 12px; text-align: left; }}
th {{ background: var(--code-bg); font-weight: 600; }}
img {{ max-width: 100%; border-radius: 6px; }}
a {{ color: var(--link); text-decoration: none; }}
a:hover {{ text-decoration: underline; }}
hr {{ border: none; border-top: 2px solid var(--border); margin: 2em 0; }}
ul,ol {{ padding-left: 2em; }}
li {{ margin: 0.25em 0; }}
input[type="checkbox"] {{ margin-right: 6px; accent-color: var(--link); }}
.katex-display {{ margin: 1em 0; overflow-x: auto; }}
</style>
<script src="/vendor/marked/marked.min.js"></script>
<script src="/vendor/katex/katex.min.js"></script>
<script src="/vendor/katex/auto-render.min.js"></script>
<script src="/vendor/highlightjs/highlight.min.js"></script>
</head><body>
<div id="content"></div>
<script>
// Configure marked
marked.setOptions({{
  gfm: true,
  breaks: true,
  highlight: function(code, lang) {{
    if (lang && hljs.getLanguage(lang)) {{
      try {{ return hljs.highlight(code, {{ language: lang }}).value; }} catch(e) {{}}
    }}
    try {{ return hljs.highlightAuto(code).value; }} catch(e) {{}}
    return code;
  }}
}});

// Render markdown
document.getElementById('content').innerHTML = marked.parse({md_json});

// Render math with KaTeX
renderMathInElement(document.getElementById('content'), {{
  delimiters: [
    {{left: "$$", right: "$$", display: true}},
    {{left: "$", right: "$", display: false}},
    {{left: "\\\\(", right: "\\\\)", display: false}},
    {{left: "\\\\[", right: "\\\\]", display: true}}
  ],
  throwOnError: false
}});
</script>
</body></html>'''
            return Response(html_content, mimetype='text/html; charset=utf-8')
        except Exception as e:
            return jsonify({'error': f'Markdown render error: {e}'}), 500

    # For binary file types (images, PDF), use send_file
    if mime_type.startswith('image/') or mime_type == 'application/pdf':
        return send_file(target, mimetype=mime_type)

    # For text-based files, serve with proper encoding
    try:
        with open(target, 'rb') as f:
            raw = f.read()
        # Try to decode as text
        for enc in ['utf-8', 'utf-8-sig', 'gbk', 'latin-1']:
            try:
                content = raw.decode(enc)
                break
            except (UnicodeDecodeError, LookupError):
                continue
        else:
            return Response(raw, mimetype=mime_type)

        # For HTML files: inject <base> tag so relative CSS/JS paths resolve correctly
        # Set base to the full file path (not just directory) so that:
        #   - Relative paths like "style.css" resolve to /preview/<dir>/style.css
        #   - Anchor links like "#section" resolve to /preview/<dir>/index.html#section
        #   (If base were just /preview/<dir>/, #links would load the directory, not the file)
        if ext in ('.html', '.htm'):
            base_href = f'/preview/{path}'
            # Inject <base> tag right after <head> or at the start of the document
            if '<head>' in content:
                content = content.replace('<head>', f'<head><base href="{base_href}">', 1)
            elif '<HEAD>' in content:
                content = content.replace('<HEAD>', f'<HEAD><base href="{base_href}">', 1)
            elif '<html>' in content:
                content = content.replace('<html>', f'<html><head><base href="{base_href}"></head>', 1)
            else:
                # No <head> tag at all — prepend it
                content = f'<head><base href="{base_href}"></head>' + content

        return Response(content, mimetype=mime_type)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@bp.route('/preview/<path:subpath>', methods=['GET'])
@handle_error
def serve_preview_file(subpath):
    """Serve static files for browser preview with correct relative path resolution.
    
    When an HTML file is previewed with a <base href="/preview/project_dir/"> tag,
    relative paths like "style.css" will resolve to /preview/project_dir/style.css.
    This route serves those files from the workspace.
    """
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    # Security: prevent directory traversal
    target = os.path.realpath(os.path.join(base, subpath))
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.isfile(target):
        return jsonify({'error': 'File not found'}), 404

    # Determine MIME type
    ext = os.path.splitext(target)[1].lower()
    mime_type = _PREVIEW_MIME_TYPES.get(ext, 'application/octet-stream')

    # For binary file types, use send_file
    if mime_type.startswith('image/') or mime_type == 'application/pdf':
        return send_file(target, mimetype=mime_type)

    # For text-based files, serve with proper encoding
    try:
        with open(target, 'rb') as f:
            raw = f.read()
        content = None
        for enc in ['utf-8', 'utf-8-sig', 'gbk', 'latin-1']:
            try:
                content = raw.decode(enc)
                break
            except (UnicodeDecodeError, LookupError):
                continue
        if content is None:
            return Response(raw, mimetype=mime_type)

        # For HTML files: inject <base> tag so relative CSS/JS paths resolve correctly
        # Set base to the full file path (not just directory) so that:
        #   - Relative paths like "style.css" resolve to /preview/<dir>/style.css
        #   - Anchor links like "#section" resolve to /preview/<dir>/index.html#section
        if ext in ('.html', '.htm'):
            base_href = f'/preview/{subpath}'
            if '<head>' in content:
                content = content.replace('<head>', f'<head><base href="{base_href}">', 1)
            elif '<HEAD>' in content:
                content = content.replace('<HEAD>', f'<HEAD><base href="{base_href}">', 1)
            elif '<html>' in content:
                content = content.replace('<html>', f'<html><head><base href="{base_href}"></head>', 1)
            else:
                content = f'<head><base href="{base_href}"></head>' + content

        return Response(content, mimetype=mime_type)
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@bp.route('/api/files/save', methods=['POST'])
@handle_error
def save_file():
    data = request.json
    path = data.get('path', '')
    content = data.get('content', '')
    create = data.get('create', False)
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    target = os.path.realpath(os.path.join(base, path))
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.exists(target) and not create:
        # Auto-create file if it doesn't exist (IDE behavior)
        os.makedirs(os.path.dirname(target), exist_ok=True)

    os.makedirs(os.path.dirname(target), exist_ok=True)
    with open(target, 'w', encoding='utf-8') as f:
        f.write(content)

    return jsonify({'ok': True, 'path': path, 'saved_at': datetime.now().isoformat()})


@bp.route('/api/files/create', methods=['POST'])
@handle_error
def create_file():
    data = request.json or {}
    path = data.get('path', '')
    is_dir = data.get('is_dir', False) or data.get('type', '') == 'directory'
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    target = os.path.realpath(os.path.join(base, path))
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if is_dir:
        os.makedirs(target, exist_ok=True)
    else:
        os.makedirs(os.path.dirname(target), exist_ok=True)
        if not os.path.exists(target):
            Path(target).touch()

    return jsonify({'ok': True, 'path': path})


@bp.route('/api/files/delete', methods=['POST'])
@handle_error
def delete_file():
    import shutil

    data = request.json
    path = data.get('path', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    target = os.path.realpath(os.path.join(base, path))
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.exists(target):
        return jsonify({'error': 'Not found'}), 404

    if os.path.isdir(target):
        shutil.rmtree(target)
    else:
        os.remove(target)

    return jsonify({'ok': True})


@bp.route('/api/files/rename', methods=['POST'])
@handle_error
def rename_file():
    data = request.json
    old_path = data.get('old_path', '')
    new_path = data.get('new_path', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    old_target = os.path.realpath(os.path.join(base, old_path))
    new_target = os.path.realpath(os.path.join(base, new_path))

    if not old_target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403
    if not new_target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.exists(old_target):
        return jsonify({'error': 'Source not found'}), 404

    os.makedirs(os.path.dirname(new_target), exist_ok=True)
    os.rename(old_target, new_target)

    return jsonify({'ok': True})


@bp.route('/api/files/open_folder', methods=['POST'])
@handle_error
def open_folder():
    data = request.json
    path = data.get('path', WORKSPACE)
    if path and os.path.isdir(path):
        config = load_config()
        config['workspace'] = path
        save_config(config)
        return jsonify({'ok': True, 'workspace': path})
    return jsonify({'error': 'Invalid folder path'}), 400


# ==================== Workspace Root Selection ====================

@bp.route('/api/workspace/info', methods=['GET'])
@handle_error
def workspace_info():
    """Get current workspace information."""
    config = load_config()
    ws = config.get('workspace', WORKSPACE)
    exists = os.path.isdir(ws)
    if not exists:
        try:
            os.makedirs(ws, exist_ok=True)
            exists = True
        except OSError:
            pass
    return jsonify({
        'workspace': ws,
        'exists': exists,
        'is_default': ws == WORKSPACE,
    })


@bp.route('/api/workspace/browse', methods=['GET'])
@handle_error
def workspace_browse():
    """Browse directories for workspace selection.
    Unlike /api/files/list, this is NOT restricted to the current workspace —
    it allows navigating the whole filesystem to pick a root directory.
    Only directories are listed (no files)."""
    path = request.args.get('path', '/')
    path = os.path.realpath(path)

    if not os.path.isdir(path):
        return jsonify({'error': 'Directory not found'}, 404)

    folders = []
    try:
        for entry in sorted(os.listdir(path)):
            full = os.path.join(path, entry)
            if os.path.isdir(full) and not entry.startswith('.'):
                folders.append({
                    'name': entry,
                    'path': full,
                })
    except PermissionError:
        return jsonify({'error': 'Permission denied'}, 403)

    return jsonify({
        'folders': folders,
        'current_path': path,
        'can_go_up': path != '/',
    })


@bp.route('/api/workspace/set', methods=['POST'])
@handle_error
def workspace_set():
    """Set the workspace directory and persist it in config."""
    data = request.json
    path = data.get('path', '')
    if not path or not os.path.isdir(path):
        return jsonify({'error': 'Invalid directory path'}), 400

    config = load_config()
    config['workspace'] = path
    save_config(config)
    return jsonify({'ok': True, 'workspace': path})


# ==================== Project Management ====================

def get_project_path():
    """Get the current project path (relative to workspace) or None."""
    config = load_config()
    return config.get('project', None)


def get_effective_base():
    """Get the effective base directory for file operations.
    When a project is open, returns the project directory.
    Otherwise returns the workspace root."""
    config = load_config()
    base = config.get('workspace', WORKSPACE)
    project = config.get('project', None)
    if project:
        project_dir = os.path.join(base, project)
        if os.path.isdir(project_dir):
            return project_dir
    return base


@bp.route('/api/project/info', methods=['GET'])
@handle_error
def project_info():
    """Get current project information."""
    config = load_config()
    project = config.get('project', None)
    base = config.get('workspace', WORKSPACE)
    if project:
        project_dir = os.path.join(base, project)
        if os.path.isdir(project_dir):
            return jsonify({
                'project': project,
                'name': os.path.basename(project),
                'path': project_dir,
                'has_git': os.path.exists(os.path.join(project_dir, '.git')),
            })
    return jsonify({'project': None, 'name': None, 'path': None})


@bp.route('/api/project/open', methods=['POST'])
@handle_error
def project_open():
    """Open a project by setting its directory as the project root."""
    data = request.json
    project_rel = data.get('project', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    if not project_rel:
        return jsonify({'error': 'Project path required'}), 400

    # Security: must be under workspace
    target = os.path.realpath(os.path.join(base, project_rel))
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.isdir(target):
        return jsonify({'error': 'Directory not found'}), 404

    config['project'] = project_rel
    save_config(config)

    return jsonify({
        'ok': True,
        'project': project_rel,
        'name': os.path.basename(project_rel),
    })


@bp.route('/api/project/create', methods=['POST'])
@handle_error
def project_create():
    """Create a new project folder in the workspace and return its relative path."""
    data = request.json
    name = data.get('name', '').strip()

    if not name:
        return jsonify({'error': '项目名称不能为空'}), 400

    # Validate name: no path separators, no leading dots
    if '/' in name or '\\' in name:
        return jsonify({'error': '项目名称不能包含路径分隔符'}), 400
    if name.startswith('.'):
        return jsonify({'error': '项目名称不能以点号开头'}), 400

    config = load_config()
    base = config.get('workspace', WORKSPACE)

    if not base or not os.path.isdir(base):
        return jsonify({'error': '工作目录未设置或不存在，请先设置工作目录'}), 400

    target = os.path.realpath(os.path.join(base, name))

    # Security: must be under workspace
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    # Check if already exists
    if os.path.exists(target):
        return jsonify({'error': f'文件夹已存在: {name}'}), 409

    try:
        os.makedirs(target, exist_ok=True)
    except OSError as e:
        return jsonify({'error': f'创建文件夹失败: {e}'}), 500

    project_rel = name
    return jsonify({
        'ok': True,
        'project': project_rel,
        'name': name,
        'path': target,
    })


@bp.route('/api/project/close', methods=['POST'])
@handle_error
def project_close():
    """Close the current project, returning to workspace view."""
    config = load_config()
    config['project'] = None
    # Clear venv_path when closing project to prevent cross-project contamination.
    # When a new project is opened, autoActivateVenv() will re-detect the correct venv.
    config['venv_path'] = ''
    save_config(config)
    return jsonify({'ok': True})


@bp.route('/api/project/list_folders', methods=['GET'])
@handle_error
def project_list_folders():
    """List folders in the workspace root for the project picker."""
    config = load_config()
    base = config.get('workspace', WORKSPACE)
    path = request.args.get('path', '')

    if path:
        target = os.path.realpath(os.path.join(base, path))
    else:
        target = os.path.realpath(base)

    # Security: must be under workspace
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.isdir(target):
        return jsonify({'folders': []})

    folders = []
    try:
        for entry in sorted(os.listdir(target)):
            full = os.path.join(target, entry)
            if os.path.isdir(full) and not entry.startswith('.'):
                rel = os.path.relpath(full, base).replace(os.sep, '/')
                has_git = os.path.exists(os.path.join(full, '.git'))
                folders.append({
                    'name': entry,
                    'path': rel,
                    'has_git': has_git,
                })
    except PermissionError:
        pass

    return jsonify({
        'folders': folders,
        'current_path': os.path.relpath(target, base).replace(os.sep, '/'),
    })


@bp.route('/api/search', methods=['POST'])
@handle_error
def search_files():
    data = request.json
    query = data.get('query', '')
    pattern = data.get('pattern', '')
    file_pattern = data.get('file_pattern', '*')
    case_sensitive = data.get('case_sensitive', False)
    use_regex = data.get('use_regex', False)
    max_results = data.get('max_results', 500)
    search_path = data.get('path', '')  # optional: limit search to a subdirectory (e.g. project dir)
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    # Determine the search root: if a project is open, default to project dir
    project = config.get('project', None)
    if search_path:
        search_root = os.path.realpath(os.path.join(base, search_path))
    elif project:
        search_root = os.path.realpath(os.path.join(base, project))
    else:
        search_root = os.path.realpath(base)

    # Security: must be under workspace
    real_base = os.path.realpath(base)
    if not search_root.startswith(real_base):
        search_root = real_base
    if not os.path.isdir(search_root):
        return jsonify({'results': [], 'total': 0})

    results = []
    search_text = pattern if pattern else query

    try:
        flags = 0 if case_sensitive else re.IGNORECASE
        if use_regex:
            regex = re.compile(search_text, flags)
        else:
            regex = re.compile(re.escape(search_text), flags)

        for root, dirs, files in os.walk(search_root):
            # Skip common ignore dirs
            dirs[:] = [d for d in dirs if d not in {'.git', '__pycache__', 'node_modules', '.venv', 'venv', '.idea', '.vscode'}]
            if len(results) >= max_results:
                break

            for fname in files:
                if len(results) >= max_results:
                    break
                # Filter by file pattern
                if file_pattern != '*' and not fnmatch.fnmatch(fname, file_pattern):
                    continue
                fpath = os.path.join(root, fname)
                try:
                    with open(fpath, 'r', encoding='utf-8', errors='ignore') as f:
                        for i, line in enumerate(f, 1):
                            if regex.search(line):
                                rel = os.path.relpath(fpath, real_base).replace(os.sep, '/')
                                results.append({
                                    'file': rel,
                                    'line': i,
                                    'col': line.lower().find(search_text.lower()) if not case_sensitive else line.find(search_text),
                                    'text': line.rstrip()[:500],
                                    'match': regex.search(line).group() if regex.search(line) else '',
                                })
                                if len(results) >= max_results:
                                    break
                except (PermissionError, OSError):
                    continue
    except re.error as e:
        return jsonify({'error': f'Invalid regex: {str(e)}'}), 400

    return jsonify({'results': results, 'total': len(results)})


@bp.route('/api/search/replace', methods=['POST'])
@handle_error
def replace_in_files():
    data = request.json
    search = data.get('search', '')
    replace = data.get('replace', '')
    file_path = data.get('file_path', '')
    case_sensitive = data.get('case_sensitive', False)
    use_regex = data.get('use_regex', False)
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    if not search:
        return jsonify({'error': 'Search text required'}), 400

    real_base = os.path.realpath(base)
    target = os.path.realpath(os.path.join(base, file_path))

    if not target.startswith(real_base):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.isfile(target):
        return jsonify({'error': 'File not found'}), 404

    try:
        with open(target, 'r', encoding='utf-8') as f:
            content = f.read()

        flags = 0 if case_sensitive else re.IGNORECASE
        if use_regex:
            new_content = re.sub(search, replace, content, flags=flags)
        else:
            new_content = re.sub(re.escape(search), replace.replace('\\', '\\\\'), content, flags=flags)

        if new_content == content:
            return jsonify({'ok': True, 'replacements': 0})

        with open(target, 'w', encoding='utf-8') as f:
            f.write(new_content)

        return jsonify({'ok': True, 'replacements': len(re.findall(search if use_regex else re.escape(search), content, flags=flags))})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ==================== File Download & Project Archive ====================

import tempfile
import zipfile
import time as _time

@bp.route('/api/files/download', methods=['GET'])
@handle_error
def download_file():
    """Download a single file or directory as zip."""
    path = request.args.get('path', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    target = os.path.realpath(os.path.join(base, path))
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.exists(target):
        return jsonify({'error': 'Path not found'}), 404

    # If it's a file, serve directly
    if os.path.isfile(target):
        ext = os.path.splitext(target)[1].lower()
        mime = _PREVIEW_MIME_TYPES.get(ext, 'application/octet-stream')
        return send_file(target, mimetype=mime, as_attachment=True,
                         download_name=os.path.basename(target))

    # If it's a directory, create a zip on-the-fly
    dirname = os.path.basename(target) or 'project'
    zip_name = f'{dirname}.zip'

    # Create temp zip
    tmp_fd, tmp_path = tempfile.mkstemp(suffix='.zip', prefix='phoneide_dl_')
    try:
        with os.fdopen(tmp_fd, 'wb') as zf:
            with zipfile.ZipFile(zf, 'w', zipfile.ZIP_DEFLATED) as z:
                for root, dirs, files in os.walk(target):
                    # Skip hidden dirs and __pycache__
                    dirs[:] = [d for d in dirs if not d.startswith('.') and d != '__pycache__' and d != 'node_modules' and d != '.git']
                    for f in files:
                        if f.startswith('.') and f != '.env':
                            continue
                        fp = os.path.join(root, f)
                        arcname = os.path.relpath(fp, target)
                        # Skip files > 50MB
                        if os.path.getsize(fp) > 50 * 1024 * 1024:
                            continue
                        z.write(fp, arcname)
        return send_file(tmp_path, mimetype='application/zip', as_attachment=True,
                         download_name=zip_name)
    except Exception as e:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return jsonify({'error': str(e)}), 500


@bp.route('/api/files/project-archive', methods=['POST'])
@handle_error
def create_project_archive():
    """Create a zip archive of the project and return a download URL.
    Used by the AI agent's project_download tool."""
    data = request.json or {}
    path = data.get('path', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)
    project = config.get('project', None)

    # Determine target directory
    if path:
        target = os.path.realpath(os.path.join(base, path))
    elif project:
        target = os.path.realpath(os.path.join(base, project))
    else:
        target = os.path.realpath(base)

    if not target.startswith(os.path.realpath(base)):
        return jsonify({'ok': False, 'error': 'Access denied'}), 403

    if not os.path.isdir(target):
        return jsonify({'ok': False, 'error': 'Target is not a directory'}), 400

    dirname = os.path.basename(target) or 'project'
    timestamp = _time.strftime('%Y%m%d_%H%M%S')
    zip_filename = f'{dirname}_{timestamp}.zip'

    # Store archives in a temp directory
    archive_dir = os.path.join(base, '.phoneide_archives')
    os.makedirs(archive_dir, exist_ok=True)
    zip_path = os.path.join(archive_dir, zip_filename)

    file_count = 0
    total_size = 0
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as z:
        for root, dirs, files in os.walk(target):
            dirs[:] = [d for d in dirs if not d.startswith('.') and d != '__pycache__' and d != 'node_modules' and d != '.git' and d != '.phoneide_archives']
            for f in files:
                if f.startswith('.') and f != '.env':
                    continue
                fp = os.path.join(root, f)
                arcname = os.path.relpath(fp, target)
                if os.path.getsize(fp) > 50 * 1024 * 1024:
                    continue
                z.write(fp, arcname)
                file_count += 1
                total_size += os.path.getsize(fp)

    zip_size = os.path.getsize(zip_path)
    # Return a download URL relative to the workspace
    rel_path = os.path.relpath(zip_path, base).replace(os.sep, '/')

    return jsonify({
        'ok': True,
        'filename': zip_filename,
        'path': rel_path,
        'download_url': f'/api/files/download?path={rel_path}',
        'file_count': file_count,
        'total_size': total_size,
        'zip_size': zip_size,
    })


@bp.route('/api/files/office-preview', methods=['GET'])
@handle_error
def office_preview():
    """Convert Office/PDF documents (docx/pptx/xlsx/pdf) to HTML for in-browser preview."""
    path = request.args.get('path', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    target = os.path.realpath(os.path.join(base, path))
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.isfile(target):
        return jsonify({'error': 'File not found'}), 404

    ext = os.path.splitext(target)[1].lower()

    try:
        if ext == '.docx':
            return _docx_to_html(target)
        elif ext == '.pptx':
            return _pptx_to_html(target)
        elif ext == '.xlsx':
            return _xlsx_to_html(target)
        elif ext == '.pdf':
            return _pdf_to_html(target)
        else:
            return jsonify({'error': f'Unsupported file type: {ext}'}), 400
    except Exception as e:
        return jsonify({'error': f'Preview failed: {str(e)}'}), 500


@bp.route('/api/files/pdf-preview', methods=['GET'])
@handle_error
def pdf_preview():
    """Render a PDF file in-browser using PDF.js for a rich reading experience.
    Provides page navigation, zoom, and text extraction — much better than
    raw send_file which relies on the browser's built-in PDF viewer (unreliable in iframes).
    Uses local PDF.js files first, falls back to CDN if local files are unavailable."""
    path = request.args.get('path', '')
    config = load_config()
    base = config.get('workspace', WORKSPACE)

    target = os.path.realpath(os.path.join(base, path))
    if not target.startswith(os.path.realpath(base)):
        return jsonify({'error': 'Access denied'}), 403

    if not os.path.isfile(target):
        return jsonify({'error': 'File not found'}), 404

    ext = os.path.splitext(target)[1].lower()
    if ext != '.pdf':
        return jsonify({'error': 'Not a PDF file'}), 400

    # Build a PDF.js-based viewer page that loads the raw PDF via /api/files/preview
    rel_path = os.path.relpath(target, os.path.realpath(base)).replace(os.sep, '/')
    from urllib.parse import quote as _url_quote
    pdf_url = f'/api/files/preview?path={_url_quote(rel_path)}'
    file_name = _html_escape(os.path.basename(target))

    # Check if local PDF.js files exist
    pdfjs_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'static', 'vendor', 'pdfjs')
    local_pdfjs_available = (os.path.isfile(os.path.join(pdfjs_dir, 'pdf.min.mjs')) and
                             os.path.isfile(os.path.join(pdfjs_dir, 'pdf.worker.min.mjs')))

    # Build import URLs: local first, CDN as fallback
    if local_pdfjs_available:
        pdfjs_main_url = '/static/vendor/pdfjs/pdf.min.mjs'
        pdfjs_worker_url = '/static/vendor/pdfjs/pdf.worker.min.mjs'
        pdfjs_fallback_main = 'https://cdn.jsdelivr.net/npm/pdfjs-dist@4.0.379/build/pdf.min.mjs'
        pdfjs_fallback_worker = 'https://cdn.jsdelivr.net/npm/pdfjs-dist@4.0.379/build/pdf.worker.min.mjs'
    else:
        pdfjs_main_url = 'https://cdn.jsdelivr.net/npm/pdfjs-dist@4.0.379/build/pdf.min.mjs'
        pdfjs_worker_url = 'https://cdn.jsdelivr.net/npm/pdfjs-dist@4.0.379/build/pdf.worker.min.mjs'
        pdfjs_fallback_main = ''
        pdfjs_fallback_worker = ''

    html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{file_name} - PDF Preview</title>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
       background: #525659; color: #333; overflow: hidden; display: flex; flex-direction: column; height: 100vh; }}
.toolbar {{ background: #323639; padding: 8px 12px; display: flex; align-items: center; gap: 10px;
            flex-shrink: 0; border-bottom: 1px solid #1e1e1e; }}
.toolbar button {{ background: #4a4d50; color: #e0e0e0; border: none; border-radius: 4px;
                   padding: 6px 12px; cursor: pointer; font-size: 13px; }}
.toolbar button:hover {{ background: #5a5d60; }}
.toolbar button:disabled {{ opacity: 0.4; cursor: default; }}
.toolbar span {{ color: #ccc; font-size: 13px; min-width: 80px; text-align: center; }}
.toolbar select {{ background: #4a4d50; color: #e0e0e0; border: none; border-radius: 4px;
                   padding: 4px 8px; font-size: 13px; }}
.toolbar .file-name {{ color: #aaa; font-size: 12px; margin-left: auto; max-width: 200px;
                       overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }}
.viewer {{ flex: 1; overflow: auto; padding: 10px; display: flex; flex-direction: column; align-items: center; }}
.page-container {{ margin: 5px 0; background: white; box-shadow: 0 1px 4px rgba(0,0,0,0.3); }}
.page-container canvas {{ display: block; }}
.loading {{ color: #aaa; padding: 40px; text-align: center; font-size: 16px; }}
.text-layer {{ position: absolute; left: 0; top: 0; right: 0; bottom: 0; overflow: hidden;
               opacity: 0.25; line-height: 1.0; }}
.text-layer > span {{ color: transparent; position: absolute; white-space: pre; cursor: text; }}
</style>
</head><body>
<div class="toolbar">
  <button id="btnPrev" title="Previous Page">&#9664; Prev</button>
  <span id="pageInfo">- / -</span>
  <button id="btnNext" title="Next Page">Next &#9654;</button>
  <select id="zoomSelect" title="Zoom Level">
    <option value="0.5">50%</option>
    <option value="0.75">75%</option>
    <option value="1" selected>100%</option>
    <option value="1.25">125%</option>
    <option value="1.5">150%</option>
    <option value="2">200%</option>
  </select>
  <button id="btnText" title="Extract Text">Extract Text</button>
  <span class="file-name">{file_name}</span>
</div>
<div class="viewer" id="viewer">
  <div class="loading" id="loadingMsg">Loading PDF...</div>
</div>
<script type="module">
const pdfUrl = "{pdf_url}";
const pdfjsMainUrl = "{pdfjs_main_url}";
const pdfjsWorkerUrl = "{pdfjs_worker_url}";
const pdfjsFallbackMain = "{pdfjs_fallback_main}";
const pdfjsFallbackWorker = "{pdfjs_fallback_worker}";

let pdfDoc = null;
let currentPage = 1;
let scale = 1.0;
let textMode = false;
let rendering = false;

const viewer = document.getElementById('viewer');
const pageInfo = document.getElementById('pageInfo');
const btnPrev = document.getElementById('btnPrev');
const btnNext = document.getElementById('btnNext');
const zoomSelect = document.getElementById('zoomSelect');
const btnText = document.getElementById('btnText');

async function loadPdfJs(mainUrl, workerUrl) {{
  const pdfjsLib = await import(mainUrl);
  pdfjsLib.GlobalWorkerOptions.workerSrc = workerUrl;
  return pdfjsLib;
}}

async function loadPDF() {{
  let pdfjsLib = null;
  // Try primary source first
  try {{
    pdfjsLib = await loadPdfJs(pdfjsMainUrl, pdfjsWorkerUrl);
  }} catch(e) {{
    console.warn('Local PDF.js failed, trying CDN fallback...', e);
    // Try CDN fallback if available
    if (pdfjsFallbackMain) {{
      try {{
        pdfjsLib = await loadPdfJs(pdfjsFallbackMain, pdfjsFallbackWorker);
      }} catch(e2) {{
        document.getElementById('loadingMsg').textContent =
          'Failed to load PDF.js library. Please check your network connection.';
        console.error('Both local and CDN PDF.js failed:', e2);
        return;
      }}
    }} else {{
      document.getElementById('loadingMsg').textContent = 'Failed to load PDF.js: ' + e.message;
      console.error(e);
      return;
    }}
  }}
  try {{
    pdfDoc = await pdfjsLib.getDocument(pdfUrl).promise;
    document.getElementById('loadingMsg').style.display = 'none';
    updatePageInfo();
    renderAllPages();
  }} catch(e) {{
    document.getElementById('loadingMsg').textContent = 'Failed to load PDF: ' + e.message;
    console.error(e);
  }}
}}

async function renderAllPages() {{
  if (rendering) return;
  rendering = true;
  viewer.innerHTML = '';
  try {{
    for (let i = 1; i <= pdfDoc.numPages; i++) {{
      const page = await pdfDoc.getPage(i);
      const viewport = page.getViewport({{ scale }});
      const container = document.createElement('div');
      container.className = 'page-container';
      container.id = 'page-' + i;
      container.style.position = 'relative';
      if (textMode) {{
        const textContent = await page.getTextContent();
        const div = document.createElement('div');
        div.style.cssText = 'padding:20px;max-width:900px;font-size:14px;line-height:1.8;white-space:pre-wrap;';
        let lastY = null;
        let lines = [];
        let currentLine = '';
        for (const item of textContent.items) {{
          if (lastY !== null && Math.abs(item.transform[5] - lastY) > 5) {{
            lines.push(currentLine);
            currentLine = '';
          }}
          currentLine += item.str;
          lastY = item.transform[5];
        }}
        if (currentLine) lines.push(currentLine);
        div.textContent = lines.join('\\n');
        container.appendChild(div);
      }} else {{
        const canvas = document.createElement('canvas');
        canvas.width = viewport.width;
        canvas.height = viewport.height;
        container.style.width = viewport.width + 'px';
        container.style.height = viewport.height + 'px';
        container.appendChild(canvas);
        const ctx = canvas.getContext('2d');
        await page.render({{ canvasContext: ctx, viewport }}).promise;
      }}
      viewer.appendChild(container);
    }}
  }} catch(e) {{
    console.error('Render error:', e);
  }}
  rendering = false;
  // Scroll to current page
  const target = document.getElementById('page-' + currentPage);
  if (target) target.scrollIntoView({{ behavior: 'auto' }});
  // Re-observe pages
  observePages();
}}

function updatePageInfo() {{
  if (!pdfDoc) return;
  pageInfo.textContent = currentPage + ' / ' + pdfDoc.numPages;
  btnPrev.disabled = currentPage <= 1;
  btnNext.disabled = currentPage >= pdfDoc.numPages;
}}

btnPrev.addEventListener('click', () => {{
  if (currentPage > 1) {{ currentPage--; updatePageInfo(); scrollToPage(); }}
}});
btnNext.addEventListener('click', () => {{
  if (pdfDoc && currentPage < pdfDoc.numPages) {{ currentPage++; updatePageInfo(); scrollToPage(); }}
}});
zoomSelect.addEventListener('change', (e) => {{
  scale = parseFloat(e.target.value);
  if (pdfDoc) renderAllPages();
}});
btnText.addEventListener('click', () => {{
  textMode = !textMode;
  btnText.style.background = textMode ? '#0078d4' : '#4a4d50';
  if (pdfDoc) renderAllPages();
}});
function scrollToPage() {{
  const target = document.getElementById('page-' + currentPage);
  if (target) target.scrollIntoView({{ behavior: 'smooth' }});
}}
// Intersection observer to update current page number on scroll
const observer = new IntersectionObserver((entries) => {{
  for (const entry of entries) {{
    if (entry.isIntersecting && entry.intersectionRatio > 0.3) {{
      const pageNum = parseInt(entry.target.id.replace('page-', ''));
      if (!isNaN(pageNum)) {{ currentPage = pageNum; updatePageInfo(); }}
    }}
  }}
}}, {{ root: viewer, threshold: 0.3 }});
function observePages() {{
  if (!pdfDoc) return;
  for (let i = 1; i <= pdfDoc.numPages; i++) {{
    const el = document.getElementById('page-' + i);
    if (el) observer.observe(el);
  }}
}}
loadPDF();
</script>
</body></html>'''
    return Response(html, mimetype='text/html; charset=utf-8')


def _pdf_to_html(filepath):
    """Convert PDF to HTML for preview (fallback when PDF.js is unavailable).
    Uses pypdf/PyPDF2 for text extraction with position-aware formatting,
    and optionally pdf2image for page image rendering."""
    import base64
    import io

    # Try pypdf first (newer), fall back to PyPDF2
    try:
        from pypdf import PdfReader
    except ImportError:
        from PyPDF2 import PdfReader

    reader = PdfReader(filepath)
    pages_html = []
    total_pages = len(reader.pages)

    # Try to use pdf2image for high-quality page rendering
    pdf2image_available = False
    try:
        from pdf2image import convert_from_path
        pdf2image_available = True
    except ImportError:
        pass

    # Page size info for aspect ratio
    page_width_pt = 612  # default letter size
    page_height_pt = 792
    try:
        first_page = reader.pages[0]
        mediabox = first_page.mediabox
        page_width_pt = float(mediabox.width)
        page_height_pt = float(mediabox.height)
    except Exception:
        pass

    for i, page in enumerate(reader.pages):
        # Method 1: Render page as image using pdf2image (best quality)
        if pdf2image_available:
            try:
                images = convert_from_path(filepath, first_page=i + 1, last_page=i + 1, dpi=150)
                if images:
                    buf = io.BytesIO()
                    images[0].save(buf, format='PNG', optimize=True)
                    b64 = base64.b64encode(buf.getvalue()).decode('ascii')
                    aspect = page_height_pt / page_width_pt * 100 if page_width_pt > 0 else 130
                    pages_html.append(f'''<div class="page">
<div class="page-header">Page {i + 1} of {total_pages}</div>
<div class="page-image" style="position:relative;padding-bottom:{aspect:.1f}%;">
<img src="data:image/png;base64,{b64}" style="position:absolute;top:0;left:0;width:100%;height:auto;" alt="Page {i + 1}">
</div>
</div>''')
                    continue
            except Exception:
                pass

        # Method 2: Extract text with position-aware formatting
        try:
            text = page.extract_text() or ''
        except Exception:
            text = ''

        if text.strip():
            text = _html_escape(text)
            formatted_text = text.replace('\n', '<br/>\n')
        else:
            # Check if page has images (can't render them in fallback mode)
            has_images = False
            try:
                if hasattr(page, 'images') and len(page.images) > 0:
                    has_images = True
                elif hasattr(page, '/Resources') and hasattr(page['/Resources'], 'get'):
                    xobject = page['/Resources'].get('/XObject', {})
                    if xobject:
                        has_images = True
            except Exception:
                pass

            if has_images:
                formatted_text = '<p style="color:#999;font-style:italic">(This page contains images that cannot be rendered in fallback mode. Please use the full viewer for best results.)</p>'
            else:
                formatted_text = '<p style="color:#999;font-style:italic">(This page has no extractable text.)</p>'

        pages_html.append(f'''<div class="page">
<div class="page-header">Page {i + 1} of {total_pages}</div>
<div class="page-content">{formatted_text}</div>
</div>''')

    body = '\n'.join(pages_html)

    # Get the download link
    config = load_config()
    base = config.get('workspace', WORKSPACE)
    try:
        rel_path = os.path.relpath(filepath, os.path.realpath(base)).replace(os.sep, '/')
    except Exception:
        rel_path = ''

    mode_label = "Image Mode (pdf2image)" if pdf2image_available else "Text Extraction Mode (fallback)"

    html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<style>
body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
       margin: 0; padding: 20px; background: #525659; color: #333; }}
.page {{ background: white; margin: 20px auto; max-width: 800px; padding: 40px 50px;
        border-radius: 2px; box-shadow: 0 2px 8px rgba(0,0,0,0.3); min-height: 200px; }}
.page-header {{ color: #999; font-size: 11px; text-align: center; margin-bottom: 20px;
               border-bottom: 1px solid #eee; padding-bottom: 8px; }}
.page-content {{ font-size: 13px; line-height: 1.8; white-space: pre-wrap; word-wrap: break-word; }}
.page-image {{ background: #f5f5f5; }}
.page-image img {{ display: block; }}
.info-bar {{ background: #323639; color: #aaa; padding: 8px 16px; font-size: 12px;
             text-align: center; position: sticky; top: 0; z-index: 50; }}
.info-bar a {{ color: #6db3f2; text-decoration: none; }}
.info-bar a:hover {{ text-decoration: underline; }}
.info-bar .mode {{ color: #888; font-style: italic; margin-left: 12px; }}
</style></head><body>
<div class="info-bar">
  PDF Document &mdash; {total_pages} page{"s" if total_pages != 1 else ""} &mdash;
  <a href="/api/files/preview?path={_html_escape(rel_path)}">Download Original PDF</a>
  <span class="mode">{mode_label}</span>
</div>
{body}
</body></html>'''
    return Response(html, mimetype='text/html; charset=utf-8')


def _docx_to_html(filepath):
    """Convert DOCX to HTML for preview with rich formatting support."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches, Cm, Emu
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    import json as _json
    import base64
    import io

    doc = Document(filepath)
    sections = []

    # --- Helper: convert docx color to CSS ---
    def _docx_color_to_css(color):
        if color is None:
            return None
        try:
            if color.rgb and str(color.rgb) != '00000000':
                return f'#{color.rgb}'
        except Exception:
            pass
        try:
            if color.theme_color is not None:
                theme_map = {
                    0: '#000000', 1: '#FFFFFF', 2: '#44546A', 3: '#E7E6e6',
                    4: '#4472C4', 5: '#ED7D31', 6: '#A5A5A5', 7: '#FFC000',
                    8: '#5B9BD5', 9: '#70AD47',
                }
                return theme_map.get(color.theme_color, None)
        except Exception:
            pass
        return None

    # --- Helper: extract image from docx as base64 data URI ---
    _image_cache = {}
    def _get_docx_image_base64(inline_or_shape):
        try:
            blip = inline_or_shape._inline.graphic.graphicData.pic.blipFill.blip
            rId = blip.embed
            if rId in _image_cache:
                return _image_cache[rId]
            rel = doc.part.rels[rId]
            image_part = rel.target_part
            image_bytes = image_part.blob
            content_type = image_part.content_type or 'image/png'
            b64 = base64.b64encode(image_bytes).decode('ascii')
            data_uri = f'data:{content_type};base64,{b64}'
            _image_cache[rId] = data_uri
            return data_uri
        except Exception:
            return None

    # --- Helper: render a paragraph with rich formatting ---
    def _render_paragraph(para, tag='p'):
        para_style_parts = []
        para_class = ''

        # Paragraph alignment
        try:
            align = para.alignment
            align_map = {
                WD_ALIGN_PARAGRAPH.LEFT: 'left',
                WD_ALIGN_PARAGRAPH.CENTER: 'center',
                WD_ALIGN_PARAGRAPH.RIGHT: 'right',
                WD_ALIGN_PARAGRAPH.JUSTIFY: 'justify',
            }
            if align in align_map:
                para_style_parts.append(f'text-align:{align_map[align]}')
        except Exception:
            pass

        # Paragraph indentation
        try:
            fmt = para.paragraph_format
            if fmt.left_indent:
                li_pt = int(fmt.left_indent) / 12700
                para_style_parts.append(f'margin-left:{li_pt:.1f}pt')
            if fmt.right_indent:
                ri_pt = int(fmt.right_indent) / 12700
                para_style_parts.append(f'margin-right:{ri_pt:.1f}pt')
            if fmt.first_line_indent:
                fi_pt = int(fmt.first_line_indent) / 12700
                para_style_parts.append(f'text-indent:{fi_pt:.1f}pt')
            if fmt.space_before:
                sb_pt = int(fmt.space_before) / 12700
                para_style_parts.append(f'margin-top:{sb_pt:.1f}pt')
            if fmt.space_after:
                sa_pt = int(fmt.space_after) / 12700
                para_style_parts.append(f'margin-bottom:{sa_pt:.1f}pt')
            if fmt.line_spacing:
                if isinstance(fmt.line_spacing, float):
                    para_style_parts.append(f'line-height:{fmt.line_spacing:.2f}')
                else:
                    ls_pt = int(fmt.line_spacing) / 12700
                    para_style_parts.append(f'line-height:{ls_pt:.1f}pt')
        except Exception:
            pass

        # Runs with rich formatting
        runs_html = []
        for run in para.runs:
            rtext = _html_escape(run.text)
            span_styles = []

            # Font color
            try:
                font = run.font
                color_css = _docx_color_to_css(font.color)
                if color_css:
                    span_styles.append(f'color:{color_css}')
            except Exception:
                pass

            # Font size
            try:
                if run.font.size:
                    sz_pt = int(run.font.size) / 12700
                    span_styles.append(f'font-size:{sz_pt:.1f}pt')
            except Exception:
                pass

            # Font name
            try:
                if run.font.name:
                    span_styles.append(f'font-family:"{run.font.name}",sans-serif')
            except Exception:
                pass

            # Highlight
            try:
                if run.font.highlight_color:
                    hl_map = {
                        1: '#FFFF00', 2: '#00FF00', 3: '#00FFFF', 4: '#FF00FF',
                        5: '#0000FF', 6: '#FF0000', 7: '#000080', 8: '#008080',
                        9: '#808080', 10: '#C0C0C0', 11: '#800080', 12: '#800000',
                        13: '#008000', 14: '#808000', 15: '#46D5DB',
                    }
                    hl = hl_map.get(run.font.highlight_color)
                    if hl:
                        span_styles.append(f'background:{hl}')
            except Exception:
                pass

            # Bold
            if run.font.bold:
                rtext = f'<strong>{rtext}</strong>'
            # Italic
            if run.font.italic:
                rtext = f'<em>{rtext}</em>'
            # Underline
            try:
                if run.font.underline:
                    rtext = f'<u>{rtext}</u>'
            except Exception:
                pass
            # Strikethrough
            try:
                if run.font.strike:
                    rtext = f'<s>{rtext}</s>'
            except Exception:
                pass
            # Superscript / Subscript
            try:
                if run.font.superscript:
                    rtext = f'<sup>{rtext}</sup>'
                elif run.font.subscript:
                    rtext = f'<sub>{rtext}</sub>'
            except Exception:
                pass

            if span_styles:
                style_str = ';'.join(span_styles)
                rtext = f'<span style="{style_str}">{rtext}</span>'
            runs_html.append(rtext)

        content = ''.join(runs_html) if runs_html else _html_escape(para.text)

        # Check for inline images (drawing elements in runs)
        try:
            from docx.oxml.ns import qn
            for run in para.runs:
                drawings = run._element.findall(qn('w:drawing'))
                for drawing in drawings:
                    blip = drawing.find('.//' + qn('a:blip'))
                    if blip is not None:
                        rId = blip.get(qn('r:embed'))
                        if rId:
                            data_uri = _image_cache.get(rId)
                            if not data_uri:
                                try:
                                    rel = doc.part.rels[rId]
                                    image_part = rel.target_part
                                    image_bytes = image_part.blob
                                    content_type = image_part.content_type or 'image/png'
                                    b64 = base64.b64encode(image_bytes).decode('ascii')
                                    data_uri = f'data:{content_type};base64,{b64}'
                                    _image_cache[rId] = data_uri
                                except Exception:
                                    data_uri = None
                            if data_uri:
                                # Try to get image dimensions
                                try:
                                    extent = drawing.find('.//' + qn('wp:extent'))
                                    if extent is not None:
                                        cx = int(extent.get('cx', 0)) / 12700
                                        cy = int(extent.get('cy', 0)) / 12700
                                        img_style = f'max-width:{cx:.0f}pt;max-height:{cy:.0f}pt'
                                    else:
                                        img_style = 'max-width:100%'
                                except Exception:
                                    img_style = 'max-width:100%'
                                runs_html.append(f'<img src="{data_uri}" style="{img_style};height:auto;margin:4px 0" />')
                                content = ''.join(runs_html)
        except Exception:
            pass

        pstyle = f' style="{";".join(para_style_parts)}"' if para_style_parts else ''
        return f'<{tag}{pstyle}>{content}</{tag}>'

    # --- Helper: render a table cell ---
    def _render_table_cell(cell, tag='td'):
        cell_style_parts = []
        # Cell background
        try:
            from docx.oxml.ns import qn
            tc = cell._element
            tcPr = tc.find(qn('w:tcPr'))
            if tcPr is not None:
                shd = tcPr.find(qn('w:shd'))
                if shd is not None:
                    fill_val = shd.get(qn('w:fill'))
                    if fill_val and fill_val.upper() != 'AUTO' and fill_val != '00000000':
                        cell_style_parts.append(f'background:#{fill_val}')
        except Exception:
            pass

        # Cell width
        try:
            if cell.width:
                w_pt = int(cell.width) / 12700
                cell_style_parts.append(f'min-width:{w_pt:.0f}pt')
        except Exception:
            pass

        # Vertical alignment
        try:
            from docx.oxml.ns import qn
            tc = cell._element
            tcPr = tc.find(qn('w:tcPr'))
            if tcPr is not None:
                vAlign = tcPr.find(qn('w:vAlign'))
                if vAlign is not None:
                    val = vAlign.get(qn('w:val'), '')
                    va_map = {'top': 'top', 'center': 'middle', 'bottom': 'bottom'}
                    if val in va_map:
                        cell_style_parts.append(f'vertical-align:{va_map[val]}')
        except Exception:
            pass

        # Render cell paragraphs
        cell_paragraphs = []
        for para in cell.paragraphs:
            cell_paragraphs.append(_render_paragraph(para, 'p'))
        cell_content = ''.join(cell_paragraphs) if cell_paragraphs else _html_escape(cell.text)

        cstyle = f' style="{";".join(cell_style_parts)}"' if cell_style_parts else ''
        return f'<{tag}{cstyle}>{cell_content}</{tag}>'

    # --- Main rendering ---
    # Use iter_inner_content to handle interleaved paragraphs and tables
    from docx.oxml.ns import qn
    body_el = doc.element.body

    for child in body_el:
        tag_name = child.tag.split('}')[-1] if '}' in child.tag else child.tag

        if tag_name == 'p':
            # Find matching paragraph object
            para = None
            for p in doc.paragraphs:
                if p._element is child:
                    para = p
                    break
            if para is None:
                continue

            style_name = para.style.name if para.style else ''

            # Check if it's a list item
            is_list = False
            try:
                pPr = child.find(qn('w:pPr'))
                if pPr is not None:
                    numPr = pPr.find(qn('w:numPr'))
                    if numPr is not None:
                        is_list = True
            except Exception:
                pass

            if is_list:
                sections.append(f'<div class="list-item">{_render_paragraph(para, "p")}</div>')
            elif 'Heading 1' in style_name:
                sections.append(_render_paragraph(para, 'h1'))
            elif 'Heading 2' in style_name:
                sections.append(_render_paragraph(para, 'h2'))
            elif 'Heading 3' in style_name:
                sections.append(_render_paragraph(para, 'h3'))
            elif 'Heading 4' in style_name:
                sections.append(_render_paragraph(para, 'h4'))
            else:
                text = para.text.strip()
                if not text and not para.runs:
                    sections.append('<br/>')
                else:
                    sections.append(_render_paragraph(para, 'p'))

        elif tag_name == 'tbl':
            # Find matching table object
            table = None
            for t in doc.tables:
                if t._element is child:
                    table = t
                    break
            if table is None:
                continue

            table_style_parts = ['border-collapse:collapse', 'width:100%']
            rows_html = []
            for ri, row in enumerate(table.rows):
                cells_html = []
                for ci, cell in enumerate(row.cells):
                    tag = 'td'  # Don't assume first row is header
                    cells_html.append(_render_table_cell(cell, tag))
                rows_html.append(f'<tr>{"".join(cells_html)}</tr>')
            table_style_str = ';'.join(table_style_parts)
            sections.append(f'<table style="{table_style_str}">{"".join(rows_html)}</table>')

    body = '\n'.join(sections)
    html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<style>
body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
       max-width: 900px; margin: 0 auto; padding: 20px; line-height: 1.6; color: #333;
       background: #fff; }}
h1 {{ font-size: 1.8em; margin-top: 1.5em; margin-bottom: 0.5em; }}
h2 {{ font-size: 1.5em; margin-top: 1.3em; margin-bottom: 0.4em; }}
h3 {{ font-size: 1.25em; margin-top: 1.2em; margin-bottom: 0.3em; }}
h4 {{ font-size: 1.1em; margin-top: 1em; margin-bottom: 0.3em; }}
table {{ margin: 1em 0; }}
th,td {{ border: 1px solid #ccc; padding: 8px 12px; text-align: left; }}
img {{ max-width: 100%; height: auto; }}
.list-item {{ margin-left: 1.5em; position: relative; }}
.list-item::before {{ content: "•"; position: absolute; left: -1.2em; }}
</style></head><body>{body}</body></html>'''
    return Response(html, mimetype='text/html; charset=utf-8')


def _pptx_to_html(filepath):
    """Convert PPTX to HTML for preview with backgrounds, shapes, images, and styling."""
    import base64
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
    from lxml import etree

    prs = Presentation(filepath)
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height

    # EMU to px conversion (1 inch = 914400 EMU, assume 96 DPI)
    def emu_to_px(emu):
        if emu is None:
            return 0
        return round(int(emu) / 914400 * 96)

    slide_w_px = emu_to_px(slide_width_emu)
    slide_h_px = emu_to_px(slide_height_emu)

    # ---- Theme color extraction ----
    # Default Office theme colors (fallback)
    _SCHEME_MAP = {
        'dk1': '#000000', 'lt1': '#FFFFFF', 'dk2': '#44546A', 'lt2': '#E7E6E6',
        'accent1': '#4472C4', 'accent2': '#ED7D31', 'accent3': '#A5A5A5',
        'accent4': '#FFC000', 'accent5': '#5B9BD5', 'accent6': '#70AD47',
        'hlink': '#0563C1', 'folHlink': '#954F72',
    }
    # Theme index to key name mapping (python-pptx uses 0-based index)
    _THEME_KEYS = ['dk1', 'lt1', 'dk2', 'lt2', 'accent1', 'accent2',
                   'accent3', 'accent4', 'accent5', 'accent6']

    # Extract actual theme colors from the PPTX file
    _pptx_theme_colors = {}
    try:
        from pptx.oxml.ns import qn
        for sm in prs.slide_masters:
            try:
                theme_el = None
                # Try to find theme through relationships
                for rel in sm.part.rels.values():
                    if 'theme' in rel.reltype:
                        theme_part = rel.target_part
                        theme_el = etree.fromstring(theme_part.blob)
                        break
                if theme_el is not None:
                    clrScheme = theme_el.find('.//' + qn('a:clrScheme'))
                    if clrScheme is not None:
                        for child in clrScheme:
                            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                            for sub in child:
                                sub_tag = sub.tag.split('}')[-1] if '}' in sub.tag else sub.tag
                                if sub_tag == 'srgbClr':
                                    val = sub.get('val', '')
                                    if val:
                                        _pptx_theme_colors[tag] = f'#{val}'
                                elif sub_tag == 'sysClr':
                                    val = sub.get('lastClr', sub.get('val', ''))
                                    if val:
                                        _pptx_theme_colors[tag] = f'#{val}'
                        break  # Use first master's theme
            except Exception:
                continue
    except Exception:
        pass

    def _resolve_scheme_color(val):
        """Resolve a scheme color name (e.g. 'dk1', 'accent1') to CSS hex color."""
        if val in _pptx_theme_colors:
            return _pptx_theme_colors[val]
        return _SCHEME_MAP.get(val, None)

    def _apply_lum_mods(color_hex, el):
        """Apply luminance modifications (tint/shade/lumMod/lumOff) to a color."""
        from pptx.oxml.ns import qn
        if el is None:
            return color_hex
        tint = el.find(qn('a:tint'))
        shade = el.find(qn('a:shade'))
        lumMod = el.find(qn('a:lumMod'))
        lumOff = el.find(qn('a:lumOff'))
        if tint is None and shade is None and lumMod is None and lumOff is None:
            return color_hex
        try:
            base = int(color_hex[1:], 16)
            r, g, b = (base >> 16) & 0xFF, (base >> 8) & 0xFF, base & 0xFF
            # Convert to HSL for more accurate tint/shade
            r1, g1, b1 = r / 255.0, g / 255.0, b / 255.0
            cmax, cmin = max(r1, g1, b1), min(r1, g1, b1)
            l = (cmax + cmin) / 2.0
            if lumMod is not None:
                factor = int(lumMod.get('val', '100000')) / 100000
                l = l * factor
            if lumOff is not None:
                offset = int(lumOff.get('val', '0')) / 100000
                l = l + offset
            if tint is not None:
                # tint = mix with white; val/100000 = amount of white
                factor = int(tint.get('val', '100000')) / 100000
                r = int(r + (255 - r) * factor)
                g = int(g + (255 - g) * factor)
                b = int(b + (255 - b) * factor)
                return f'#{r:02X}{g:02X}{b:02X}'
            if shade is not None:
                # shade = mix with black; val/100000 = amount of black kept
                factor = int(shade.get('val', '100000')) / 100000
                r = int(r * factor)
                g = int(g * factor)
                b = int(b * factor)
                return f'#{r:02X}{g:02X}{b:02X}'
            # Fallback: simple lumMod/lumOff via RGB scaling
            r = min(255, max(0, int(r1 * 255)))
            g = min(255, max(0, int(g1 * 255)))
            b = min(255, max(0, int(b1 * 255)))
            return f'#{r:02X}{g:02X}{b:02X}'
        except Exception:
            return color_hex

    def _xml_color_to_css(el):
        """Convert an XML color element (srgbClr, schemeClr, prstClr) to CSS hex color."""
        if el is None:
            return None
        from pptx.oxml.ns import qn
        tag = el.tag.split('}')[-1] if '}' in el.tag else el.tag
        if tag == 'srgbClr':
            val = el.get('val', '000000')
            return _apply_lum_mods(f'#{val}', el)
        elif tag == 'schemeClr':
            val = el.get('val', '')
            color = _resolve_scheme_color(val)
            if color:
                return _apply_lum_mods(color, el)
            return None
        elif tag == 'prstClr':
            name = el.get('val', 'black').lower()
            prst_map = {
                'black': '#000000', 'white': '#FFFFFF', 'red': '#FF0000',
                'green': '#008000', 'blue': '#0000FF', 'yellow': '#FFFF00',
                'cyan': '#00FFFF', 'magenta': '#FF00FF',
            }
            color = prst_map.get(name)
            if color:
                return _apply_lum_mods(color, el)
            return None
        return None

    def _get_fill_css_from_xml(fill_el, default_color=None):
        """Get CSS background from a fill XML element (solidFill, gradFill, pattFill)."""
        from pptx.oxml.ns import qn
        if fill_el is None:
            return default_color or 'transparent'
        tag = fill_el.tag.split('}')[-1] if '}' in fill_el.tag else fill_el.tag
        if tag == 'solidFill':
            for child in fill_el:
                color = _xml_color_to_css(child)
                if color:
                    return color
        elif tag == 'gradFill':
            gsLst = fill_el.find(qn('a:gsLst'))
            if gsLst is not None:
                stops = []
                for gs in gsLst.findall(qn('a:gs')):
                    pos = int(gs.get('pos', '0')) / 1000
                    color = None
                    for child in gs:
                        color = _xml_color_to_css(child)
                        if color:
                            break
                    if color:
                        stops.append(f'{color} {pos}%')
                if stops:
                    lin = fill_el.find(qn('a:lin'))
                    if lin is not None:
                        ang = int(lin.get('ang', '0'))
                        css_ang = (ang / 60000) % 360
                        direction = f'{css_ang:.0f}deg'
                    else:
                        direction = 'to bottom'
                    return f'linear-gradient({direction}, {", ".join(stops)})'
        elif tag == 'pattFill':
            fgClr = fill_el.find(qn('a:fgClr'))
            if fgClr is not None:
                for child in fgClr:
                    color = _xml_color_to_css(child)
                    if color:
                        return color
        return default_color or 'transparent'

    def _get_fill_first_color(fill_el):
        """Get the first solid color from a fill XML element (for text color use)."""
        from pptx.oxml.ns import qn
        if fill_el is None:
            return None
        tag = fill_el.tag.split('}')[-1] if '}' in fill_el.tag else fill_el.tag
        if tag == 'solidFill':
            for child in fill_el:
                color = _xml_color_to_css(child)
                if color:
                    return color
        elif tag == 'gradFill':
            # For text color, just use the first gradient stop color
            gsLst = fill_el.find(qn('a:gsLst'))
            if gsLst is not None:
                for gs in gsLst.findall(qn('a:gs')):
                    for child in gs:
                        color = _xml_color_to_css(child)
                        if color:
                            return color
        return None

    def _get_fill_css(fill, default_color=None):
        """Get CSS background from a python-pptx fill object."""
        try:
            if fill.type is None:
                return default_color or 'transparent'
            if fill.type == 1:  # SOLID
                try:
                    color_css = _color_to_css(fill.fore_color)
                    if color_css:
                        return color_css
                except Exception:
                    pass
                # Try XML fallback
                try:
                    fill_el = fill._fill
                    if fill_el is not None:
                        css = _get_fill_css_from_xml(fill_el)
                        if css and css != 'transparent':
                            return css
                except Exception:
                    pass
                return default_color or 'transparent'
            elif fill.type == 2:  # GRADIENT
                try:
                    fill_el = fill._fill
                    if fill_el is not None:
                        css = _get_fill_css_from_xml(fill_el, default_color)
                        if css and css != 'transparent':
                            return css
                except Exception:
                    pass
                try:
                    color_css = _color_to_css(fill.fore_color)
                    if color_css:
                        return color_css
                except Exception:
                    pass
                return default_color or 'transparent'
            elif fill.type == 3:  # PATTERN
                try:
                    color_css = _color_to_css(fill.fore_color)
                    if color_css:
                        return color_css
                except Exception:
                    pass
                return default_color or 'transparent'
            elif fill.type == 5:  # BACKGROUND (inherit)
                return 'transparent'
        except Exception:
            pass
        return default_color or 'transparent'

    def _color_to_css(color):
        """Convert a python-pptx color object to CSS hex color string."""
        if color is None:
            return None
        try:
            # Try RGB first
            if hasattr(color, 'rgb') and color.rgb:
                rgb_str = str(color.rgb)
                if rgb_str and rgb_str != '000000ZZ':
                    return f'#{rgb_str}'
        except Exception:
            pass
        try:
            # Try theme color
            if hasattr(color, 'theme') and color.theme is not None:
                idx = color.theme
                if idx < len(_THEME_KEYS):
                    key = _THEME_KEYS[idx]
                    resolved = _resolve_scheme_color(key)
                    if resolved:
                        # Try to apply brightness from color object
                        try:
                            brightness = color.brightness
                            if brightness:
                                # brightness is a float 0-1; mix with white
                                base = int(resolved[1:], 16)
                                r, g, b = (base >> 16) & 0xFF, (base >> 8) & 0xFF, base & 0xFF
                                r = int(r + (255 - r) * brightness)
                                g = int(g + (255 - g) * brightness)
                                b = int(b + (255 - b) * brightness)
                                return f'#{r:02X}{g:02X}{b:02X}'
                        except Exception:
                            pass
                        return resolved
        except Exception:
            pass
        return None

    def _get_bg_from_xml_element(bg_el):
        """Get CSS background from a p:bg XML element."""
        from pptx.oxml.ns import qn
        if bg_el is None:
            return None
        # bgRef — reference to theme
        bgRef = bg_el.find(qn('p:bgRef'))
        if bgRef is not None:
            idx = int(bgRef.get('idx', '0'))
            if idx < 1000 and idx < len(_THEME_KEYS):
                resolved = _resolve_scheme_color(_THEME_KEYS[idx])
                if resolved:
                    return resolved
            return '#FFFFFF'
        # bgPr — background properties
        bgPr = bg_el.find(qn('p:bgPr'))
        if bgPr is not None:
            for fill_tag in ['a:solidFill', 'a:gradFill', 'a:pattFill']:
                fill_el = bgPr.find(qn(fill_tag))
                if fill_el is not None:
                    css = _get_fill_css_from_xml(fill_el)
                    if css and css != 'transparent':
                        return css
        return None

    def _get_bg_image_css(bg_el, part):
        """Get CSS for background image if present."""
        from pptx.oxml.ns import qn
        if bg_el is None:
            return None, None
        bgPr = bg_el.find(qn('p:bgPr'))
        if bgPr is None:
            return None, None
        blipFill = bgPr.find(qn('a:blipFill'))
        if blipFill is None:
            return None, None
        blip = blipFill.find(qn('a:blip'))
        if blip is None:
            return None, None
        rId = blip.get(qn('r:embed'))
        if not rId:
            return None, None
        data_uri = _get_image_base64(part, rId)
        if not data_uri:
            return None, None
        stretch = blipFill.find(qn('a:stretch'))
        tile = blipFill.find(qn('a:tile'))
        if stretch is not None:
            return f'url("{data_uri}")', 'center/cover no-repeat'
        elif tile is not None:
            return f'url("{data_uri}")', 'repeat'
        else:
            return f'url("{data_uri}")', 'center/contain no-repeat'

    def _get_slide_bg_css(slide):
        """Get CSS background for a slide, checking slide -> layout -> master chain."""
        from pptx.oxml.ns import qn

        # 1. Check slide's own background via API
        try:
            fill = slide.background.fill
            bg_css = _get_fill_css(fill)
            if bg_css and bg_css != 'transparent':
                return bg_css, None, None
        except Exception:
            pass

        # 2. Check slide's own background via XML
        cSld = slide._element.find(qn('p:cSld'))
        if cSld is not None:
            bg_el = cSld.find(qn('p:bg'))
            if bg_el is not None:
                img_css, img_size = _get_bg_image_css(bg_el, slide.part)
                if img_css:
                    return img_css, img_size, None
                bg_css = _get_bg_from_xml_element(bg_el)
                if bg_css and bg_css != 'transparent':
                    return bg_css, None, None

        # 3. Check slide layout background
        try:
            layout = slide.slide_layout
            if layout is not None:
                cSld_layout = layout._element.find(qn('p:cSld'))
                if cSld_layout is not None:
                    bg_el = cSld_layout.find(qn('p:bg'))
                    if bg_el is not None:
                        img_css, img_size = _get_bg_image_css(bg_el, slide.part)
                        if img_css:
                            return img_css, img_size, None
                        bg_css = _get_bg_from_xml_element(bg_el)
                        if bg_css and bg_css != 'transparent':
                            return bg_css, None, None
        except Exception:
            pass

        # 4. Check slide master background
        try:
            master = slide.slide_layout.slide_master
            if master is not None:
                cSld_master = master._element.find(qn('p:cSld'))
                if cSld_master is not None:
                    bg_el = cSld_master.find(qn('p:bg'))
                    if bg_el is not None:
                        img_css, img_size = _get_bg_image_css(bg_el, slide.part)
                        if img_css:
                            return img_css, img_size, None
                        bg_css = _get_bg_from_xml_element(bg_el)
                        if bg_css and bg_css != 'transparent':
                            return bg_css, None, None
        except Exception:
            pass

        return '#FFFFFF', None, None  # Default white

    def _get_image_base64(part, rId):
        """Extract image from PPTX by relationship ID and return base64 data URI."""
        try:
            rel = part.rels[rId]
            image_part = rel.target_part
            image_bytes = image_part.blob
            content_type = image_part.content_type
            if not content_type:
                ext = image_part.partname.split('.')[-1].lower()
                ct_map = {'png': 'image/png', 'jpg': 'image/jpeg', 'jpeg': 'image/jpeg',
                          'gif': 'image/gif', 'bmp': 'image/bmp', 'svg': 'image/svg+xml',
                          'tiff': 'image/tiff', 'tif': 'image/tiff'}
                content_type = ct_map.get(ext, 'image/png')
            b64 = base64.b64encode(image_bytes).decode('ascii')
            return f'data:{content_type};base64,{b64}'
        except Exception:
            return None

    def _render_text_frame(shape):
        """Render text frame from a shape, reading XML for accurate formatting.
        
        Reads <a:defRPr> (paragraph default run properties) and <a:rPr> (run-level
        overrides) and merges them, since python-pptx's API only exposes run-level
        properties, missing the inherited defaults.
        """
        from pptx.oxml.ns import qn

        # Try XML-based rendering first for accuracy
        sp = shape._element
        txBody = sp.find(qn('p:txBody'))
        if txBody is None:
            txBody = sp.find(qn('a:txBody'))
        if txBody is None:
            return _render_text_frame_api(shape)

        paragraphs_html = []
        for p_el in txBody.findall(qn('a:p')):
            pPr = p_el.find(qn('a:pPr'))
            para_style = []

            # Paragraph alignment
            if pPr is not None:
                algn = pPr.get('algn')
                algn_map = {'l': 'left', 'ctr': 'center', 'r': 'right', 'just': 'justify', 'dist': 'justify'}
                if algn in algn_map:
                    para_style.append(f'text-align:{algn_map[algn]}')

            # Paragraph spacing
            if pPr is not None:
                spcAft = pPr.find(qn('a:spcAft'))
                if spcAft is not None:
                    spcPts = spcAft.find(qn('a:spcPts'))
                    if spcPts is not None:
                        val = int(spcPts.get('val', '0')) / 100
                        para_style.append(f'margin-bottom:{val:.0f}pt')
                spcBef = pPr.find(qn('a:spcBef'))
                if spcBef is not None:
                    spcPts = spcBef.find(qn('a:spcPts'))
                    if spcPts is not None:
                        val = int(spcPts.get('val', '0')) / 100
                        para_style.append(f'margin-top:{val:.0f}pt')

            # Default run properties (defRPr) — inherited by all runs in this paragraph
            def_rpr = {}
            if pPr is not None:
                defRPr_el = pPr.find(qn('a:defRPr'))
                if defRPr_el is not None:
                    def_rpr = _parse_rpr(defRPr_el)

            # Render runs
            runs_html = []
            r_els = p_el.findall(qn('a:r'))

            if not r_els:
                # Check for field elements (slide numbers, dates, etc.)
                fld_els = p_el.findall(qn('a:fld'))
                if fld_els:
                    for fld in fld_els:
                        t_el = fld.find(qn('a:t'))
                        if t_el is not None and t_el.text:
                            runs_html.append(_html_escape(t_el.text))
                else:
                    runs_html.append('&nbsp;')  # Empty line placeholder
            else:
                for r_el in r_els:
                    t_el = r_el.find(qn('a:t'))
                    if t_el is None:
                        continue
                    rtext = _html_escape(t_el.text or '')

                    # Merge defRPr (defaults) with rPr (run overrides)
                    run_props = dict(def_rpr)
                    rPr_el = r_el.find(qn('a:rPr'))
                    if rPr_el is not None:
                        run_overrides = _parse_rpr(rPr_el)
                        run_props.update(run_overrides)

                    # Build inline styles
                    run_styles = []

                    # Color — only set if explicitly specified
                    color = run_props.get('color')
                    if color:
                        run_styles.append(f'color:{color}')

                    # Font size
                    sz = run_props.get('font_size')
                    if sz:
                        run_styles.append(f'font-size:{sz}')

                    # Font name
                    fn = run_props.get('font_name')
                    if fn:
                        run_styles.append(f'font-family:"{fn}",sans-serif')

                    # Bold / Italic / Underline / Strikethrough
                    if run_props.get('bold'):
                        rtext = f'<strong>{rtext}</strong>'
                    if run_props.get('italic'):
                        rtext = f'<em>{rtext}</em>'
                    if run_props.get('underline'):
                        rtext = f'<u>{rtext}</u>'
                    if run_props.get('strike'):
                        rtext = f'<s>{rtext}</s>'

                    if run_styles:
                        style_str = ';'.join(run_styles)
                        rtext = f'<span style="{style_str}">{rtext}</span>'

                    runs_html.append(rtext)

            pstyle = f' style="{";".join(para_style)}"' if para_style else ''
            paragraphs_html.append(f'<p{pstyle}>{"".join(runs_html)}</p>')

        return ''.join(paragraphs_html)

    def _parse_rpr(rpr_el):
        """Parse an <a:rPr> or <a:defRPr> element into a dict of formatting properties."""
        from pptx.oxml.ns import qn
        props = {}

        # Font size (sz is in hundredths of a point, e.g. 1800 = 18pt)
        sz = rpr_el.get('sz')
        if sz:
            sz_pt = int(sz) / 100
            props['font_size'] = f'{sz_pt:.0f}pt'

        # Bold
        b = rpr_el.get('b')
        if b is not None:
            props['bold'] = b != '0'

        # Italic
        i = rpr_el.get('i')
        if i is not None:
            props['italic'] = i != '0'

        # Underline
        u = rpr_el.get('u')
        if u is not None:
            props['underline'] = u != 'none'

        # Strikethrough
        strike = rpr_el.get('strike')
        if strike is not None:
            props['strike'] = strike != 'noStrike'

        # Font name — from latin/ea/cs elements
        # Skip theme font references like +mn-lt, +mj-lt (major/minor theme fonts)
        latin = rpr_el.find(qn('a:latin'))
        if latin is not None:
            typeface = latin.get('typeface', '')
            if typeface and not typeface.startswith('+'):
                props['font_name'] = typeface
        if 'font_name' not in props:
            ea = rpr_el.find(qn('a:ea'))
            if ea is not None:
                typeface = ea.get('typeface', '')
                if typeface and not typeface.startswith('+'):
                    props['font_name'] = typeface

        # Text color — from solidFill only (gradient fill is not a valid CSS color)
        solidFill = rpr_el.find(qn('a:solidFill'))
        if solidFill is not None:
            for child in solidFill:
                color = _xml_color_to_css(child)
                if color:
                    props['color'] = color
                    break

        # NOTE: Do NOT use gradFill as text color — CSS color property doesn't
        # support gradients. If needed, extract the first stop color only.

        return props

    def _render_text_frame_api(shape):
        """Fallback: render text frame using python-pptx API."""
        text_parts = []
        for para in shape.text_frame.paragraphs:
            runs_html = []
            para_style = []
            try:
                align = para.alignment
                align_map = {1: 'left', 2: 'center', 3: 'right', 4: 'justify', 5: 'center', 6: 'left', 7: 'right'}
                if align in align_map:
                    para_style.append(f'text-align:{align_map[align]}')
            except Exception:
                pass
            for run in para.runs:
                rtext = _html_escape(run.text)
                run_styles = []
                try:
                    color_css = _color_to_css(run.font.color)
                    if color_css:
                        run_styles.append(f'color:{color_css}')
                except Exception:
                    pass
                try:
                    if run.font.size:
                        sz_pt = int(run.font.size) / 12700
                        run_styles.append(f'font-size:{sz_pt:.0f}pt')
                except Exception:
                    pass
                try:
                    if run.font.name:
                        run_styles.append(f'font-family:"{run.font.name}",sans-serif')
                except Exception:
                    pass
                if run.font.bold:
                    rtext = f'<strong>{rtext}</strong>'
                if run.font.italic:
                    rtext = f'<em>{rtext}</em>'
                try:
                    if run.font.underline:
                        rtext = f'<u>{rtext}</u>'
                except Exception:
                    pass
                if run_styles:
                    style_str = ';'.join(run_styles)
                    rtext = f'<span style="{style_str}">{rtext}</span>'
                runs_html.append(rtext)
            text = ''.join(runs_html) if runs_html else _html_escape(para.text)
            pstyle = f' style="{";".join(para_style)}"' if para_style else ''
            text_parts.append(f'<p{pstyle}>{text}</p>')
        return ''.join(text_parts)

    def _render_shape(slide, shape):
        """Render a single shape to HTML with positioning and styling."""
        from pptx.oxml.ns import qn

        # Skip shapes with no position or zero-size
        try:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            if left is None or top is None or width is None or height is None:
                return ''
        except Exception:
            return ''

        left_px = emu_to_px(left)
        top_px = emu_to_px(top)
        width_px = emu_to_px(width)
        height_px = emu_to_px(height)

        # Skip zero-size shapes
        if width_px <= 0 or height_px <= 0:
            return ''

        style_parts = [
            f'position:absolute',
            f'left:{left_px}px',
            f'top:{top_px}px',
            f'width:{width_px}px',
            f'height:{height_px}px',
            f'box-sizing:border-box',
            f'overflow:hidden',
        ]

        # ---- Shape fill ----
        fill_rendered = False
        # Try python-pptx API first
        try:
            fill = shape.fill
            fill_css = _get_fill_css(fill)
            if fill_css and fill_css != 'transparent':
                style_parts.append(f'background:{fill_css}')
                fill_rendered = True
        except Exception:
            pass

        # Try XML directly if API didn't work
        if not fill_rendered:
            try:
                sp = shape._element
                # Look for spPr (shape properties) which contains fill info
                spPr = sp.find(qn('p:spPr'))
                if spPr is None:
                    spPr = sp.find('.//' + qn('a:spPr'))
                if spPr is not None:
                    for fill_tag in ['a:solidFill', 'a:gradFill', 'a:pattFill', 'a:blipFill']:
                        fill_el = spPr.find(qn(fill_tag))
                        if fill_el is not None:
                            if fill_tag == 'a:blipFill':
                                blip = fill_el.find(qn('a:blip'))
                                if blip is not None:
                                    rId = blip.get(qn('r:embed'))
                                    if rId:
                                        data_uri = _get_image_base64(slide.part, rId)
                                        if data_uri:
                                            style_parts.append(f'background-image:url("{data_uri}")')
                                            style_parts.append('background-size:cover')
                                            style_parts.append('background-position:center')
                                            fill_rendered = True
                                            break
                            else:
                                css = _get_fill_css_from_xml(fill_el)
                                if css and css != 'transparent':
                                    style_parts.append(f'background:{css}')
                                    fill_rendered = True
                                    break
            except Exception:
                pass

        # ---- Shape border/line ----
        try:
            line = shape.line
            if line.fill.type is not None:
                if line.fill.type == 1:  # solid
                    line_color = _color_to_css(line.color)
                    line_width = line.width
                    if line_color and line_width:
                        w_pt = int(line_width) / 12700
                        style_parts.append(f'border:{w_pt:.1f}pt solid {line_color}')
                    elif line_color:
                        style_parts.append(f'border:1pt solid {line_color}')
                elif line.fill.type == 5:  # no line
                    style_parts.append('border:none')
        except Exception:
            pass

        # ---- Rotation ----
        try:
            rotation = shape.rotation
            if rotation:
                style_parts.append(f'transform:rotate({rotation}deg)')
        except Exception:
            pass

        # ---- Determine shape type and render content ----
        shape_type = None
        try:
            shape_type = shape.shape_type
        except Exception:
            pass

        content = ''

        # Picture shape
        if shape_type == MSO_SHAPE_TYPE.PICTURE or shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE:
            try:
                sp = shape._element
                blipFill_el = sp.find('.//' + qn('a:blip'))
                if blipFill_el is not None:
                    rId = blipFill_el.get(qn('r:embed'))
                    if rId:
                        data_uri = _get_image_base64(slide.part, rId)
                        if data_uri:
                            content = f'<img src="{data_uri}" style="width:100%;height:100%;object-fit:contain">'
            except Exception:
                pass

        # Group shape — recurse into children
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                group_html = []
                for child in shape.shapes:
                    group_html.append(_render_shape(slide, child))
                content = ''.join(group_html)
                style_parts.append('overflow:visible')
            except Exception:
                pass

        # Text frame — use XML-based rendering for accurate formatting
        if not content and shape.has_text_frame:
            content = _render_text_frame(shape)
            # Add padding for text shapes without background
            if not fill_rendered:
                style_parts.append('padding:4px 8px')

        # Table
        elif not content and shape.has_table:
            table = shape.table
            rows_html = []
            for ri, row in enumerate(table.rows):
                cells = []
                for cell in row.cells:
                    cell_text = _html_escape(cell.text)
                    cell_style = []
                    try:
                        cell_fill = cell.fill
                        cell_bg = _get_fill_css(cell_fill)
                        if cell_bg and cell_bg != 'transparent':
                            cell_style.append(f'background:{cell_bg}')
                    except Exception:
                        pass
                    cstyle = f' style="{";".join(cell_style)}"' if cell_style else ''
                    cells.append(f'<td{cstyle}>{cell_text}</td>')
                rows_html.append(f'<tr>{"".join(cells)}</tr>')
            content = f'<table border="1" cellpadding="6" cellspacing="0" style="width:100%;border-collapse:collapse;font-size:14px">{"".join(rows_html)}</table>'

        style_str = ';'.join(style_parts)
        return f'<div class="shape" style="{style_str}">{content}</div>'

    # ---- Build each slide ----
    slides_html = []
    for i, slide in enumerate(prs.slides):
        bg_css, bg_size, _ = _get_slide_bg_css(slide)
        bg_style = f'background:{bg_css}'
        if bg_size:
            bg_style = f'background-image:{bg_css};background-size:{bg_size}'

        # Render ONLY the slide's own shapes (not layout/master shapes)
        # Layout/master decorative shapes are often placeholders or visual elements
        # that get inherited through the background chain — rendering them separately
        # causes overlap and duplication issues.
        shapes_html = []
        for shape in slide.shapes:
            try:
                rendered = _render_shape(slide, shape)
                if rendered:
                    shapes_html.append(rendered)
            except Exception:
                # Fallback: extract text at least
                try:
                    if shape.has_text_frame:
                        text = _html_escape(shape.text)
                        shapes_html.append(
                            f'<div class="shape" style="position:absolute;'
                            f'left:{emu_to_px(shape.left)}px;top:{emu_to_px(shape.top)}px;'
                            f'width:{emu_to_px(shape.width)}px;font-size:18px">'
                            f'{text}</div>'
                        )
                except Exception:
                    pass

        shapes_content = '\n'.join(shapes_html)
        slides_html.append(f'''<div class="slide" style="{bg_style}">
<div class="slide-scaler">
<div class="slide-canvas" style="position:relative;width:{slide_w_px}px;height:{slide_h_px}px;overflow:hidden">
{shapes_content}
</div>
</div>
<div class="slide-number">Slide {i + 1} / {len(prs.slides)}</div>
</div>''')

    body = '\n'.join(slides_html)
    html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<style>
body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
       margin: 0; padding: 20px; background: #f0f0f0; color: #333; }}
.slide {{ margin: 20px auto; max-width: {slide_w_px}px; width: calc(100vw - 40px);
          border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.15); overflow: hidden; }}
.slide-scaler {{ width: 100%; position: relative; }}
.slide-canvas {{ transform-origin: top left; }}
.slide-number {{ color: rgba(0,0,0,0.5); font-size: 11px; padding: 6px 12px; text-align: center; background: #fff; }}
.shape {{ box-sizing: border-box; overflow: hidden; word-wrap: break-word; }}
.shape p {{ margin: 2px 0; line-height: 1.3; word-wrap: break-word; }}
.shape img {{ display: block; }}
table {{ border-collapse: collapse; width: 100%; margin: 4px 0; }}
th,td {{ border: 1px solid #ddd; padding: 6px 10px; text-align: left; font-size: 13px; }}
th {{ background: #f5f5f5; font-weight: bold; }}
</style></head><body>{body}
<script>
(function(){{
var sw={slide_w_px},sh={slide_h_px};
function resize(){{
document.querySelectorAll('.slide').forEach(function(s){{
var c=s.querySelector('.slide-canvas');
var sc=s.querySelector('.slide-scaler');
if(!c||!sc)return;
var w=s.clientWidth;
var scale=Math.min(1,w/sw);
c.style.transform='scale('+scale+')';
sc.style.height=(sh*scale)+'px';
}});
}}
resize();
window.addEventListener('resize',resize);
}})();
</script></body></html>'''
    return Response(html, mimetype='text/html; charset=utf-8')


def _xlsx_to_html(filepath):
    """Convert XLSX to HTML for preview with cell styling, merged cells, and formatting."""
    from openpyxl import load_workbook
    from openpyxl.styles import (
        PatternFill, Font, Border, Side, Alignment, numbers
    )
    from openpyxl.utils import get_column_letter

    wb = load_workbook(filepath, data_only=False)
    sheets_html = []

    def _xlsx_color_to_css(color):
        """Convert openpyxl Color to CSS color string."""
        if color is None:
            return None
        try:
            if color.type == 'rgb' and color.rgb and str(color.rgb) != '00000000':
                rgb = str(color.rgb)
                if len(rgb) == 8:
                    # ARGB format — skip alpha for CSS
                    return f'#{rgb[2:]}'
                elif len(rgb) == 6:
                    return f'#{rgb}'
            elif color.type == 'theme':
                # Theme colors
                theme_map = {
                    0: '#000000', 1: '#FFFFFF', 2: '#44546A', 3: '#E7E6E6',
                    4: '#4472C4', 5: '#ED7D31', 6: '#A5A5A5', 7: '#FFC000',
                    8: '#5B9BD5', 9: '#70AD47',
                }
                idx = color.theme if color.theme is not None else 0
                base = theme_map.get(idx, '#000000')
                # Apply tint
                tint = color.tint if color.tint else 0
                if tint != 0:
                    try:
                        base_int = int(base[1:], 16)
                        r, g, b = (base_int >> 16) & 0xFF, (base_int >> 8) & 0xFF, base_int & 0xFF
                        if tint > 0:
                            # Lighten
                            r = min(255, int(r + (255 - r) * tint))
                            g = min(255, int(g + (255 - g) * tint))
                            b = min(255, int(b + (255 - b) * tint))
                        else:
                            # Darken
                            factor = 1 + tint
                            r = max(0, int(r * factor))
                            g = max(0, int(g * factor))
                            b = max(0, int(b * factor))
                        return f'#{r:02X}{g:02X}{b:02X}'
                    except Exception:
                        pass
                return base
            elif color.type == 'indexed':
                # Indexed colors — simplified mapping
                idx_colors = {
                    0: '#000000', 1: '#FFFFFF', 2: '#FF0000', 3: '#00FF00',
                    4: '#0000FF', 5: '#FFFF00', 6: '#FF00FF', 7: '#00FFFF',
                    8: '#000000', 9: '#FFFFFF', 10: '#FF0000', 11: '#00FF00',
                    12: '#0000FF', 13: '#FFFF00', 14: '#FF00FF', 15: '#00FFFF',
                    16: '#800000', 17: '#008000', 18: '#000080', 19: '#808000',
                    20: '#800080', 21: '#008080', 22: '#C0C0C0', 23: '#808080',
                    24: '#9999FF', 25: '#993366', 26: '#FFFFCC', 27: '#CCFFFF',
                    28: '#660066', 29: '#FF8080', 30: '#0066CC', 31: '#CCCCFF',
                    32: '#000080', 33: '#FF00FF', 34: '#FFFF00', 35: '#00FFFF',
                    36: '#800080', 37: '#800000', 38: '#008080', 39: '#0000FF',
                    40: '#00CCFF', 41: '#CCFFFF', 42: '#CCFFCC', 43: '#FFFF99',
                    44: '#99CCFF', 45: '#FF99CC', 46: '#CC99FF', 47: '#FFCC99',
                    48: '#3366FF', 49: '#33CCCC', 50: '#99CC00', 51: '#FFCC00',
                    52: '#FF9900', 53: '#FF6600', 54: '#666699', 55: '#969696',
                    56: '#003366', 57: '#339966', 58: '#003300', 59: '#333300',
                    60: '#993300', 61: '#993366', 62: '#333399', 63: '#333333',
                }
                return idx_colors.get(color.indexed, '#000000')
        except Exception:
            pass
        return None

    def _format_cell_value(cell):
        """Format cell value according to its number format."""
        if cell.value is None:
            return ''

        # If it's a string, return as-is
        if isinstance(cell.value, str):
            return _html_escape(cell.value)

        # Try to apply number formatting
        try:
            fmt = cell.number_format
            if fmt and fmt != 'General':
                # Use openpyxl's formatting
                try:
                    formatted = numbers.format_number(cell.value, fmt)
                    return _html_escape(str(formatted))
                except Exception:
                    pass
        except Exception:
            pass

        # Default formatting for numbers
        if isinstance(cell.value, float):
            # Check if it looks like an integer
            if cell.value == int(cell.value):
                return _html_escape(str(int(cell.value)))
            return _html_escape(f'{cell.value:g}')

        return _html_escape(str(cell.value))

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]

        # Collect column widths
        col_widths = {}
        for col_idx in range(1, min(ws.max_column + 1, 51)):
            col_letter = get_column_letter(col_idx)
            if col_letter in ws.column_dimensions and ws.column_dimensions[col_letter].width:
                col_widths[col_idx] = min(ws.column_dimensions[col_letter].width * 8, 400)
            else:
                col_widths[col_idx] = 80  # Default width

        # Collect merged cell ranges
        merged_ranges = {}
        for merge_range in ws.merged_cells.ranges:
            min_row = merge_range.min_row
            min_col = merge_range.min_col
            max_row = merge_range.max_row
            max_col = merge_range.max_col
            rowspan = max_row - min_row + 1
            colspan = max_col - min_col + 1
            merged_ranges[(min_row, min_col)] = (rowspan, colspan, max_row, max_col)

        # Track cells that are part of a merge (not the top-left)
        merged_slave = set()
        for (r, c), (rs, cs, mr, mc) in merged_ranges.items():
            for rr in range(r, mr + 1):
                for cc in range(c, mc + 1):
                    if (rr, cc) != (r, c):
                        merged_slave.add((rr, cc))

        rows_html = []
        max_row = min(ws.max_row, 500)
        max_col = min(ws.max_column, 50)

        for ri, row in enumerate(ws.iter_rows(min_row=1, max_row=max_row, max_col=max_col), 1):
            cells = []
            for cell in row:
                ci = cell.column

                # Skip cells that are part of a merge (not top-left)
                if (ri, ci) in merged_slave:
                    continue

                # Format cell value
                val_str = _format_cell_value(cell)

                # Build cell style
                cell_style_parts = []

                # Cell background fill
                try:
                    fill = cell.fill
                    if fill and fill.start_color and fill.patternType and fill.patternType != 'none':
                        bg_color = _xlsx_color_to_css(fill.start_color)
                        if bg_color:
                            cell_style_parts.append(f'background:{bg_color}')
                except Exception:
                    pass

                # Font color
                try:
                    font = cell.font
                    if font.color:
                        font_color = _xlsx_color_to_css(font.color)
                        if font_color:
                            cell_style_parts.append(f'color:{font_color}')
                except Exception:
                    pass

                # Font bold
                try:
                    if cell.font.bold:
                        cell_style_parts.append('font-weight:bold')
                except Exception:
                    pass

                # Font italic
                try:
                    if cell.font.italic:
                        cell_style_parts.append('font-style:italic')
                except Exception:
                    pass

                # Font size
                try:
                    if cell.font.size:
                        cell_style_parts.append(f'font-size:{cell.font.size}pt')
                except Exception:
                    pass

                # Font name
                try:
                    if cell.font.name:
                        cell_style_parts.append(f'font-family:"{cell.font.name}",sans-serif')
                except Exception:
                    pass

                # Text underline
                try:
                    if cell.font.underline:
                        cell_style_parts.append('text-decoration:underline')
                except Exception:
                    pass

                # Text strikethrough
                try:
                    if cell.font.strikethrough:
                        cell_style_parts.append('text-decoration:line-through')
                except Exception:
                    pass

                # Cell alignment
                try:
                    alignment = cell.alignment
                    if alignment:
                        h_map = {
                            'left': 'left', 'center': 'center', 'right': 'right',
                            'fill': 'left', 'justify': 'justify', 'centerContinuous': 'center',
                            'distributed': 'justify',
                        }
                        if alignment.horizontal and alignment.horizontal in h_map:
                            cell_style_parts.append(f'text-align:{h_map[alignment.horizontal]}')
                        elif isinstance(cell.value, (int, float)):
                            # Numbers default to right align
                            cell_style_parts.append('text-align:right')

                        v_map = {'top': 'top', 'center': 'middle', 'bottom': 'bottom'}
                        if alignment.vertical and alignment.vertical in v_map:
                            cell_style_parts.append(f'vertical-align:{v_map[alignment.vertical]}')

                        if alignment.wrap_text:
                            cell_style_parts.append('white-space:normal;word-wrap:break-word')
                except Exception:
                    pass

                # Cell borders
                try:
                    border = cell.border
                    if border:
                        for side_name, side_obj in [('border-top', border.top), ('border-right', border.right),
                                                      ('border-bottom', border.bottom), ('border-left', border.left)]:
                            if side_obj and side_obj.style and side_obj.style != 'none':
                                b_color = _xlsx_color_to_css(side_obj.color) if side_obj.color else '#000'
                                b_style_map = {
                                    'thin': '1px solid', 'medium': '2px solid',
                                    'thick': '3px solid', 'dotted': '1px dotted',
                                    'dashed': '1px dashed', 'mediumDashed': '2px dashed',
                                    'hair': '1px solid', 'double': '3px double',
                                    'mediumDotDash': '2px dotted', 'slantDashDot': '2px dashed',
                                }
                                b_css = b_style_map.get(side_obj.style, '1px solid')
                                cell_style_parts.append(f'{side_name}:{b_css} {b_color}')
                except Exception:
                    pass

                # Column width (via min-width)
                if ci in col_widths:
                    cell_style_parts.append(f'min-width:{col_widths[ci]:.0f}px')

                # Merge cell attributes
                merge_attrs = ''
                if (ri, ci) in merged_ranges:
                    rowspan, colspan, _, _ = merged_ranges[(ri, ci)]
                    if rowspan > 1:
                        merge_attrs += f' rowspan="{rowspan}"'
                    if colspan > 1:
                        merge_attrs += f' colspan="{colspan}"'

                cstyle = f' style="{";".join(cell_style_parts)}"' if cell_style_parts else ''
                cells.append(f'<td{cstyle}{merge_attrs}>{val_str}</td>')

            if cells:
                rows_html.append(f'<tr>{"".join(cells)}</tr>')

        table_html = f'<table>{"".join(rows_html)}</table>'
        sheets_html.append(f'''<div class="sheet">
<h3 class="sheet-title">{_html_escape(sheet_name)}</h3>
{table_html}
</div>''')

    wb.close()
    body = '\n'.join(sheets_html)

    # If multiple sheets, add tab navigation
    if len(wb.sheetnames) > 1:
        tabs = ''.join(f'<button class="sheet-tab" onclick="showSheet({i})">{_html_escape(name)}</button>'
                       for i, name in enumerate(wb.sheetnames))
        nav = f'<div class="sheet-nav">{tabs}</div>'
    else:
        nav = ''

    html = f'''<!DOCTYPE html>
<html><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<style>
body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
       max-width: 1200px; margin: 0 auto; padding: 20px; line-height: 1.4; color: #333; }}
.sheet {{ margin: 16px 0; }}
.sheet-title {{ margin-bottom: 8px; color: #555; font-size: 16px; }}
.sheet-nav {{ display: flex; gap: 8px; margin-bottom: 16px; flex-wrap: wrap; }}
.sheet-tab {{ padding: 6px 16px; border: 1px solid #ccc; border-radius: 4px; background: #f8f8f8;
              cursor: pointer; font-size: 14px; }}
.sheet-tab.active {{ background: #0078d4; color: white; border-color: #0078d4; }}
table {{ border-collapse: collapse; width: 100%; margin: 0 0 1em; font-size: 13px; }}
td, th {{ border: 1px solid #d0d0d0; padding: 4px 8px; text-align: left; white-space: nowrap;
          overflow: hidden; text-overflow: ellipsis; }}
td {{ min-width: 50px; }}
</style></head><body>
{nav}
{body}
<script>
(function() {{
    const sheets = document.querySelectorAll('.sheet');
    const tabs = document.querySelectorAll('.sheet-tab');
    if (sheets.length > 1) {{
        sheets.forEach((s, i) => {{ if (i > 0) s.style.display = 'none'; }});
        if (tabs[0]) tabs[0].classList.add('active');
    }}
    window.showSheet = function(idx) {{
        sheets.forEach((s, i) => {{ s.style.display = i === idx ? '' : 'none'; }});
        tabs.forEach((t, i) => {{ t.classList.toggle('active', i === idx); }});
    }};
}})();
</script>
</body></html>'''
    return Response(html, mimetype='text/html; charset=utf-8')


def _html_escape(text):
    """Simple HTML entity escaping."""
    if not isinstance(text, str):
        text = str(text)
    return text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('"', '&quot;')
